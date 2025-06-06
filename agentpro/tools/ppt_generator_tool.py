from agentpro.tools.base_tool import Tool
from pptx import Presentation
from pptx.util import Pt
import os
import re

class PowerPointGeneratorTool(Tool):
    name: str = "ppt_generator"
    description: str = "Generates a PowerPoint summary from multiple academic paper summaries"
    action_type: str = "ppt_generator"
    input_format: str = "A dictionary with {'topic': str, 'summaries': list of str or dict}"

    def run(self, data: dict) -> str:
        topic = data.get("topic", "Research Summary").strip()
        summaries = data.get("summaries", [])

        if not summaries:
            return "No summaries provided for the PowerPoint."

        # Clean file name from topic
        filename_base = re.sub(r'[^\w\s-]', '', topic).strip().replace(' ', '_').lower()
        file_path = f"{filename_base}_summary.pptx"

        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        bullet_slide_layout = prs.slide_layouts[1]

        # Slide 1: Title
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = topic
        slide.placeholders[1].text = "Academic Research Summary"

        # Content Slides with paper titles
        for summary_entry in summaries:
            if isinstance(summary_entry, dict):
                paper_title = summary_entry.get("title") or topic
                summary_text = summary_entry.get("summary", "")
                authors = summary_entry.get("authors", "Unknown")
            else:
                paper_title = topic
                summary_text = summary_entry
                authors = "Unknown"

            clean_summary = self.clean_summary_text(summary_text)
            slide = prs.slides.add_slide(bullet_slide_layout)

            # Set slide title
            title_shape = slide.shapes.title
            title_shape.text = paper_title
            for paragraph in title_shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(30)

            # Set slide body
            textbox = slide.placeholders[1]
            textbox.text = f"Authors: {authors}\n\n{clean_summary}"
            self.format_textbox(textbox, 16)

        # Save the presentation
        prs.save(file_path)
        return f"PowerPoint presentation generated: {file_path}"

    def clean_summary_text(self, text: str) -> str:
        return re.sub(
            r"^(This\s(article|paper|study|work)\s(aims|discusses|examines|presents|reviews)\s)",
            "", text.strip(), flags=re.IGNORECASE)

    def format_textbox(self, textbox, font_size_pt: int = 16):
        for paragraph in textbox.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(font_size_pt)
        textbox.text_frame.word_wrap = True
